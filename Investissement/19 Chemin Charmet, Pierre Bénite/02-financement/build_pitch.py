#!/usr/bin/env python3
"""Pitch deck Pierre-Bénite — version épurée banque, gris anthracite."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pathlib import Path

# === Palette gris anthracite ===
ANTHRACITE = RGBColor(0x36, 0x45, 0x4F)
ANTHRACITE_DARK = RGBColor(0x21, 0x2A, 0x33)
GOLD = RGBColor(0xC9, 0xA9, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GREY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GREY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x00, 0x70, 0x33)
RED = RGBColor(0x9C, 0x00, 0x06)
SOFT_BG = RGBColor(0xFA, 0xFA, 0xFA)

FONT_T = 'Georgia'
FONT_B = 'Calibri'

OUT = Path(__file__).parent / "Pitch-Pierre-Benite.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def txt(slide, l, t, w, h, text, size=14, bold=False, color=DARK, font=FONT_B, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rect(slide, l, t, w, h, color=ANTHRACITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def bg(slide, color=SOFT_BG):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background()
    spTree = b._element.getparent()
    spTree.remove(b._element)
    spTree.insert(2, b._element)


def header(slide, title, sub=None):
    rect(slide, 0, 0, 0.3, 7.5, color=ANTHRACITE)
    txt(slide, 0.7, 0.4, 12, 0.7, title, size=30, bold=True, color=ANTHRACITE, font=FONT_T)
    if sub:
        txt(slide, 0.7, 1.15, 12, 0.4, sub, size=13, color=GREY, font=FONT_B)


# ============ SLIDE 1 : COUVERTURE ============
s = blank(); bg(s, ANTHRACITE)
# Trait or
rect(s, 0, 6, 13.333, 0.04, color=GOLD)
txt(s, 0.5, 2.4, 12.3, 1.2, "PIERRE-BÉNITE", size=64, bold=True, color=WHITE, font=FONT_T, align=PP_ALIGN.CENTER)
txt(s, 0.5, 3.6, 12.3, 0.5, "19 Chemin Charmet — limite Saint-Genis-Laval", size=22, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 0.5, 4.4, 12.3, 0.5, "Projet d'investissement immobilier", size=16, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 0.5, 6.4, 12.3, 0.4, "Marc Chenut  &  Antoine Bocquet", size=14, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 0.5, 6.9, 12.3, 0.3, "Mai 2026", size=10, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)


# ============ SLIDE 2 : L'OPPORTUNITÉ (4 encarts) ============
s = blank(); bg(s)
header(s, "L'opportunité en bref")

# 4 encarts : Décote / Plus-value latente / Cash flow LCD / Cash flow LLD
callouts = [
    ("DÉCOTE À L'ACHAT",
     "1 167 €/m² acheté\nvs 4 752 €/m² médiane maison Saint-Genis-Laval (DVF)\n→ -75 % vs marché",
     ANTHRACITE),
    ("PLUS-VALUE LATENTE",
     "Prêt sollicité :    588 K€\nValeur revente :   1 260 K€\n→ Marge : +672 K€",
     GOLD),
    ("CASH FLOW LCD",
     "Revenus net charges :   +10 712 €/mois\nMensualité prêt :              -3 502 €/mois\n→ Cash flow : +7 210 €/mois",
     GREEN),
    ("CASH FLOW LONGUE DURÉE",
     "Loyers (T2 + T3) :          +5 450 €/mois\nCharges + Prêt :              -4 081 €/mois\n→ Cash flow : +1 369 €/mois",
     ANTHRACITE_DARK),
]
for i, (t1, t2, col) in enumerate(callouts):
    l = 0.7 + (i % 2) * 6
    tp = 1.9 + (i // 2) * 2.6
    rect(s, l, tp, 5.7, 2.3, color=col)
    txt(s, l + 0.25, tp + 0.2, 5.3, 0.5, t1, size=14, bold=True, color=GOLD if col != GOLD else ANTHRACITE_DARK, font=FONT_B)
    txt(s, l + 0.25, tp + 0.8, 5.3, 1.4, t2, size=13, color=WHITE if col != GOLD else ANTHRACITE_DARK, font=FONT_B)


# ============ SLIDE 3 : L'ÉQUIPE ============
s = blank(); bg(s)
header(s, "L'équipe", "Profils complémentaires — opérateur LCD + expert travaux")
for i, (name, role, pts) in enumerate([
    ("Marc Chenut", "Opérateur LCD — 6 ans d'expérience", [
        "97 % d'occupation prouvée sur 4 opérations actives à Lyon",
        "Auteur du livre « La Méthode 97 % »",
        "Salarié Neocamino (3 000 € /mois net)",
        "Société M&Ms : bénéfice ~30 K€/an",
        "Trésorerie disponible : 100 K€",
    ]),
    ("Antoine Bocquet", "Directeur Rhône-Alpes — RenovationMan", [
        "Réseau d'artisans → travaux -50 % vs particulier",
        "Maîtrise délais et qualité de rénovation",
        "EURL personnelle",
        "Compagne actuaire : 7 000 €/mois net",
        "Trésorerie disponible : 100 K€",
    ]),
]):
    l = 0.7 + i * 6.3
    txt(s, l, 1.95, 5.7, 0.6, name, size=24, bold=True, color=ANTHRACITE, font=FONT_T)
    txt(s, l, 2.55, 5.7, 0.4, role, size=14, color=GOLD, font=FONT_B)
    tb = s.shapes.add_textbox(Inches(l), Inches(3.15), Inches(5.7), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, point in enumerate(pts):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "•  " + point
        for r in p.runs:
            r.font.name = FONT_B
            r.font.size = Pt(13)
            r.font.color.rgb = DARK
        p.space_after = Pt(10)


# ============ SLIDE 4 : LE BIEN ============
s = blank(); bg(s)
header(s, "Le bien", "Maison ~300 m² — limite Saint-Genis-Laval — 10 min métro Lyon")
stats = [
    ("ADRESSE", "19 Chemin Charmet, 69310 Pierre-Bénite"),
    ("SURFACE", "~300 m²"),
    ("MARCHÉ DE RÉFÉRENCE", "Saint-Genis-Laval"),
    ("ACCÈS", "10 min à pied du métro Lyon"),
    ("DIVISION", "5 lots — 3 T3 (dont 2 avec jardin) + 2 T2"),
    ("STANDING", "Rénovation haut de gamme"),
]
for i, (k, v) in enumerate(stats):
    r0 = i // 2
    c0 = i % 2
    l = 0.7 + c0 * 6.3
    tp = 2.0 + r0 * 1.55
    rect(s, l, tp, 5.7, 1.35, color=GREY_LIGHT)
    txt(s, l + 0.25, tp + 0.2, 5.3, 0.4, k, size=10, bold=True, color=ANTHRACITE, font=FONT_B)
    txt(s, l + 0.25, tp + 0.6, 5.3, 0.65, v, size=15, bold=True, color=DARK, font=FONT_B)


# ============ SLIDE 5 : LE MARCHÉ (refait propre) ============
s = blank(); bg(s)
header(s, "Le marché — Saint-Genis-Laval", "DVF / MeilleursAgents / efficity — mai 2026")

# Big stat : médiane Saint-Genis
rect(s, 0.7, 1.9, 12, 1.5, color=ANTHRACITE)
txt(s, 0.7, 2.0, 12, 0.45, "Prix m² appartements Saint-Genis-Laval", size=14, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 0.7, 2.5, 12, 0.85, "Médiane 3 480 €/m²  —  Moyenne 3 950 €/m²  —  Top quartiers jusqu'à 4 991 €/m²",
    size=20, bold=True, color=WHITE, font=FONT_T, align=PP_ALIGN.CENTER)

# Notre positionnement
rect(s, 0.7, 3.7, 5.7, 1.4, color=GOLD)
txt(s, 0.7, 3.8, 5.7, 0.4, "NOTRE ACHAT", size=12, bold=True, color=ANTHRACITE_DARK, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 0.7, 4.2, 5.7, 0.5, "1 167 €/m²", size=26, bold=True, color=ANTHRACITE_DARK, font=FONT_T, align=PP_ALIGN.CENTER)
txt(s, 0.7, 4.75, 5.7, 0.35, "soit -75 % vs médiane maison Saint-Genis", size=11, color=ANTHRACITE_DARK, font=FONT_B, align=PP_ALIGN.CENTER)

rect(s, 7, 3.7, 5.7, 1.4, color=GREEN)
txt(s, 7, 3.8, 5.7, 0.4, "NOTRE REVENTE CIBLE", size=12, bold=True, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 7, 4.2, 5.7, 0.5, "4 200 €/m²", size=26, bold=True, color=WHITE, font=FONT_T, align=PP_ALIGN.CENTER)
txt(s, 7, 4.75, 5.7, 0.35, "300 m² × 4 200 = 1 260 K€", size=11, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)

# Plus-value
rect(s, 2.5, 5.5, 8.3, 1.5, color=ANTHRACITE_DARK)
txt(s, 2.5, 5.6, 8.3, 0.4, "PLUS-VALUE LATENTE", size=12, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 2.5, 6.0, 8.3, 0.7, "+ 522 K€", size=32, bold=True, color=WHITE, font=FONT_T, align=PP_ALIGN.CENTER)
txt(s, 2.5, 6.7, 8.3, 0.3, "soit +71 % de marge sur le coût total opération", size=11, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)


# ============ SLIDE 6 : CONCURRENCE HÔTELIÈRE ============
s = blank(); bg(s)
header(s, "Marché hôtelier — Zone de chalandise", "Marché 2-3 étoiles : dormir 1-2 nuits, propre, calme, prix correct")

l = 0.7; tp = 1.9
total_w = 11.9
cols = [("Hôtel", 8), ("Catégorie", 1.9), ("Chambres", 2)]
rect(s, l, tp, total_w, 0.5, color=ANTHRACITE)
col_l = l
for name, w in cols:
    txt(s, col_l + 0.2, tp + 0.1, w - 0.2, 0.4, name, size=12, bold=True, color=WHITE, font=FONT_B, align=PP_ALIGN.LEFT)
    col_l += w

rows = [
    ("Ibis Styles Lyon Sud", "3 ★", "115"),
    ("Première Classe Irigny", "2 ★", "106"),
    ("Kyriad Lyon Sud", "3 ★", "61"),
    ("Ibis Budget Saint-Genis-Laval", "2 ★", "57"),
    ("Ibis Budget Lyon Sud — Saint-Fons A7", "2 ★", "42"),
    ("Hotel Lyon Sud", "3 ★", "40"),
    ("Le Saint Vincent", "2 ★", "32"),
    ("Campanile Prime", "3 ★", "~70"),
    ("Hotel Muse La Baguette", "2-3 ★", "~30"),
    ("TOTAL chambres concurrentes", "", "~553"),
]
rtop = tp + 0.5
for i, row in enumerate(rows):
    is_total = (i == len(rows) - 1)
    rc = GOLD if is_total else (GREY_LIGHT if i % 2 == 1 else WHITE)
    rect(s, l, rtop + i * 0.4, total_w, 0.4, color=rc)
    col_l = l
    for j, val in enumerate(row):
        w = cols[j][1]
        b = is_total or j == 0
        col = ANTHRACITE_DARK if is_total else (ANTHRACITE if j == 0 else DARK)
        txt(s, col_l + 0.2, rtop + i * 0.4 + 0.08, w - 0.2, 0.28, val, size=11, bold=b, color=col, font=FONT_B)
        col_l += w

# Argument
rect(s, 0.7, 6.6, 12, 0.7, color=ANTHRACITE_DARK)
txt(s, 0.85, 6.65, 11.7, 0.3, "Hôtel rentable dès 60 % d'occupation → ~330 chambres louées/jour sur le marché.",
    size=11, color=GOLD, font=FONT_B)
txt(s, 0.85, 6.95, 11.7, 0.3, "Notre offre : 5 lots × 97 % d'occupation = 4,85 lots/jour, soit 1,5 % du flux concurrent — pas de saturation.",
    size=11, bold=True, color=WHITE, font=FONT_B)


# ============ SLIDE 7 : PLAN DE FINANCEMENT (tableau propre) ============
s = blank(); bg(s)
header(s, "Plan de financement", "Coût total 738 K€ — apport 150 K€ — emprunt 588 K€")

l = 1; tp = 2
total_w = 11
cols_h = [("Poste", 6), ("Montant", 5)]
rect(s, l, tp, total_w, 0.6, color=ANTHRACITE)
col_l = l
for name, w in cols_h:
    txt(s, col_l + 0.25, tp + 0.15, w - 0.25, 0.4, name, size=13, bold=True, color=WHITE, font=FONT_B)
    col_l += w

rows = [
    ("Achat", "350 000 €"),
    ("Frais de notaire (8 %)", "28 000 €"),
    ("Travaux (rénovation HG)", "300 000 €"),
    ("Ameublement", "50 000 €"),
    ("Frais bancaires (dossier + caution)", "10 000 €"),
    ("COÛT TOTAL OPÉRATION", "738 000 €"),
    ("Apport (Marc 75 K + Antoine 75 K)", "150 000 €"),
    ("EMPRUNT SOLLICITÉ", "588 000 €"),
]
rtop = tp + 0.6
for i, (lab, val) in enumerate(rows):
    is_total = (lab == "COÛT TOTAL OPÉRATION" or lab == "EMPRUNT SOLLICITÉ")
    rc = GOLD if is_total else (GREY_LIGHT if i % 2 == 1 else WHITE)
    txt_col = ANTHRACITE_DARK if is_total else DARK
    rect(s, l, rtop + i * 0.5, total_w, 0.5, color=rc)
    txt(s, l + 0.25, rtop + i * 0.5 + 0.1, 5.5, 0.35, lab, size=12, bold=is_total, color=txt_col, font=FONT_B)
    txt(s, l + 6.2, rtop + i * 0.5 + 0.1, 4.5, 0.35, val, size=13, bold=is_total or i > 4, color=txt_col, font=FONT_B, align=PP_ALIGN.RIGHT)


# ============ SLIDE 8 : PLANS A & B (LCD + LLD côte à côte) ============
s = blank(); bg(s)
header(s, "Plan A — LCD  &  Plan B — Location longue durée", "Revenus mensuels — régime plein")

# Plan A LCD (left)
rect(s, 0.5, 1.9, 6, 0.6, color=GREEN)
txt(s, 0.5, 2.0, 6, 0.45, "PLAN A — Location courte durée", size=14, bold=True, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)

l = 0.5; tp = 2.5
data_a = [
    ("Revenus NET de charges", "+10 712 €"),
    ("Mensualité prêt", "-3 502 €"),
    ("Cash flow AVANT IMPÔT", "+7 210 €"),
]
for i, (lab, val) in enumerate(data_a):
    is_total = (i == 2)
    rc = GREEN if is_total else (GREY_LIGHT if i % 2 == 1 else WHITE)
    txt_col = WHITE if is_total else DARK
    rect(s, l, tp + i * 0.7, 6, 0.7, color=rc)
    txt(s, l + 0.2, tp + i * 0.7 + 0.18, 3.5, 0.35, lab, size=12, bold=is_total, color=txt_col, font=FONT_B)
    txt(s, l + 3.8, tp + i * 0.7 + 0.15, 2, 0.4, val, size=16 if is_total else 14, bold=True, color=txt_col, font=FONT_B, align=PP_ALIGN.RIGHT)

# Plan B LLD (right)
rect(s, 6.8, 1.9, 6, 0.6, color=ANTHRACITE)
txt(s, 6.8, 2.0, 6, 0.45, "PLAN B — Location longue durée", size=14, bold=True, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)

l = 6.8; tp = 2.5
data_b = [
    ("Loyer T2 (2 lots × 850 €)", "+1 700 €"),
    ("Loyer T3 (3 lots × 1 250 €)", "+3 750 €"),
    ("    Total loyers", "+5 450 €"),
    ("Charges propriétaire", "-579 €"),
    ("Mensualité prêt", "-3 502 €"),
    ("Cash flow AVANT IMPÔT", "+1 369 €"),
]
for i, (lab, val) in enumerate(data_b):
    is_subtotal = (i == 2)
    is_total = (i == 5)
    if is_total:
        rc = ANTHRACITE
        txt_col = WHITE
    elif is_subtotal:
        rc = GREY_LIGHT
        txt_col = ANTHRACITE
    else:
        rc = GREY_LIGHT if i % 2 == 1 else WHITE
        txt_col = DARK
    rect(s, l, tp + i * 0.45, 6, 0.45, color=rc)
    txt(s, l + 0.2, tp + i * 0.45 + 0.07, 3.5, 0.32, lab, size=11, bold=is_total or is_subtotal, color=txt_col, font=FONT_B)
    txt(s, l + 3.8, tp + i * 0.45 + 0.07, 2, 0.32, val, size=14 if is_total else 12, bold=True, color=txt_col, font=FONT_B, align=PP_ALIGN.RIGHT)

# Note
txt(s, 0.5, 6.7, 12.3, 0.4, "Plan B (location nue meublée valeurs marché Saint-Genis) sécurise le projet : couvre intégralement crédit et charges.",
    size=11, color=GREY, font=FONT_B, align=PP_ALIGN.CENTER)


# ============ SLIDE 9 : CASH FLOW CUMULÉ 3 ANS ============
s = blank(); bg(s)
header(s, "Cash flow projet — Cumulé sur 3 ans", "Travaux M1-M6 + Montée M7-M9 + Régime plein M10+ / Différé partiel 18 mois")

cd = CategoryChartData()
cd.categories = ['Fin An 1', 'Fin An 2', 'Fin An 3']
cd.add_series('Cash flow cumulé (K€)', (33, 129, 215))

cs = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.5), Inches(2), Inches(10), Inches(4), cd)
c = cs.chart
c.has_title = False
c.has_legend = False
c.category_axis.tick_labels.font.size = Pt(12)
c.category_axis.tick_labels.font.name = FONT_B
c.value_axis.tick_labels.font.size = Pt(10)
c.value_axis.tick_labels.font.name = FONT_B
p = c.plots[0]
p.has_data_labels = True
p.data_labels.font.size = Pt(14)
p.data_labels.font.bold = True

txt(s, 1, 6.3, 11, 0.5, "Cumulé fin An 3 : ~215 K€ de cash flow projet généré (avant impôt)",
    size=15, bold=True, color=GREEN, font=FONT_B, align=PP_ALIGN.CENTER)


# ============ SLIDE 10 : MONTAGE JURIDIQUE ============
s = blank(); bg(s)
header(s, "Structure juridique", "4 sociétés à constituer — séparation patrimoine / exploitation")

rect(s, 1.5, 2, 3.5, 1, color=ANTHRACITE)
txt(s, 1.5, 2.1, 3.5, 0.4, "Holding Marc (SAS)", size=14, bold=True, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 1.5, 2.55, 3.5, 0.4, "50 %", size=12, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)

rect(s, 8.5, 2, 3.5, 1, color=ANTHRACITE)
txt(s, 8.5, 2.1, 3.5, 0.4, "Holding Antoine (SAS)", size=14, bold=True, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 8.5, 2.55, 3.5, 0.4, "50 %", size=12, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)

ar1 = s.shapes.add_connector(1, Inches(3.25), Inches(3), Inches(5.5), Inches(3.7))
ar1.line.color.rgb = ANTHRACITE
ar1.line.width = Pt(2)
ar2 = s.shapes.add_connector(1, Inches(10.25), Inches(3), Inches(8), Inches(3.7))
ar2.line.color.rgb = ANTHRACITE
ar2.line.width = Pt(2)

rect(s, 5, 3.7, 3.5, 1, color=GOLD)
txt(s, 5, 3.8, 3.5, 0.4, "SCI à l'IS", size=18, bold=True, color=ANTHRACITE_DARK, font=FONT_T, align=PP_ALIGN.CENTER)
txt(s, 5, 4.25, 3.5, 0.4, "Détient le bien", size=11, color=ANTHRACITE_DARK, font=FONT_B, align=PP_ALIGN.CENTER)

ar3 = s.shapes.add_connector(1, Inches(6.75), Inches(4.7), Inches(6.75), Inches(5.5))
ar3.line.color.rgb = ANTHRACITE
ar3.line.width = Pt(2)
txt(s, 7, 4.9, 3, 0.4, "bail commercial", size=10, color=GOLD, font=FONT_B)

rect(s, 5, 5.5, 3.5, 1, color=ANTHRACITE)
txt(s, 5, 5.6, 3.5, 0.4, "SAS Conciergerie", size=14, bold=True, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 5, 6.05, 3.5, 0.4, "Exploite — encaisse Airbnb/Booking", size=11, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)

txt(s, 0.7, 6.85, 12, 0.4, "Séparation propriétaire / exploitant : protection patrimoine + optimisation fiscale (mère-fille)",
    size=11, color=DARK, font=FONT_B, align=PP_ALIGN.CENTER)


# ============ SLIDE 11 : STRATÉGIE FISCALE ============
s = blank(); bg(s)
header(s, "Stratégie fiscale", "IS 15 % partout — bénéfices équilibrés sous seuil 42 500 €")

for i, (name, sub, pts, color, txtcol) in enumerate([
    ("Conciergerie SAS", "Bénéfice ~42 K€/an", [
        "CA Airbnb / Booking en direct",
        "Charges variables (ménage, blanchisserie...)",
        "Loyer commercial → SCI",
        "IS 15 % : ~6 400 €/an",
    ], ANTHRACITE, WHITE),
    ("SCI à l'IS", "Bénéfice ~23 K€/an", [
        "Encaisse loyer commercial",
        "Amort. bât. 30 ans + travaux 15 ans + mobilier 10 ans",
        "Intérêts emprunt déductibles",
        "IS 15 % : ~3 400 €/an",
    ], GOLD, ANTHRACITE_DARK),
]):
    l = 0.7 + i * 6.3
    rect(s, l, 1.85, 5.7, 1, color=color)
    txt(s, l + 0.2, 1.95, 5.3, 0.5, name, size=20, bold=True, color=txtcol, font=FONT_T)
    txt(s, l + 0.2, 2.5, 5.3, 0.4, sub, size=11, color=txtcol, font=FONT_B)
    for j, pt in enumerate(pts):
        tb = s.shapes.add_textbox(Inches(l + 0.2), Inches(3.1 + j * 0.65), Inches(5.5), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "•  " + pt
        for r in p.runs:
            r.font.name = FONT_B
            r.font.size = Pt(13)
            r.font.color.rgb = DARK

txt(s, 0.7, 6.5, 12, 0.5, "Total IS : ~9 800 €/an (15 % consolidé) au lieu de 25 % si on dépassait le seuil",
    size=14, bold=True, color=GREEN, font=FONT_B, align=PP_ALIGN.CENTER)


# ============ SLIDE 12 : LA DEMANDE + PARTENARIAT ============
s = blank(); bg(s, ANTHRACITE)
txt(s, 0.5, 0.5, 12.3, 0.7, "Notre demande", size=40, bold=True, color=WHITE, font=FONT_T, align=PP_ALIGN.CENTER)
txt(s, 0.5, 1.25, 12.3, 0.4, "Conditions d'emprunt souhaitées", size=14, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)

# Top — conditions
conds = [
    ("Montant", "588 K€"),
    ("Durée", "20 ans"),
    ("Taux", "3,8 % visé"),
    ("Différé", "18 mois"),
]
for i, (k, v) in enumerate(conds):
    l = 0.7 + i * 3.05
    tp = 1.9
    rect(s, l, tp, 2.8, 1.3, color=GOLD)
    txt(s, l + 0.1, tp + 0.15, 2.6, 0.35, k.upper(), size=11, bold=True, color=ANTHRACITE_DARK, font=FONT_B, align=PP_ALIGN.CENTER)
    txt(s, l + 0.1, tp + 0.5, 2.6, 0.7, v, size=20, bold=True, color=ANTHRACITE_DARK, font=FONT_T, align=PP_ALIGN.CENTER)

# Bottom — partenariat
rect(s, 0.7, 3.6, 12, 3.5, color=ANTHRACITE_DARK)
txt(s, 0.9, 3.7, 11.6, 0.5, "Nous cherchons un partenaire bancaire long terme", size=18, bold=True, color=GOLD, font=FONT_T, align=PP_ALIGN.LEFT)

points = [
    "Plusieurs opérations prévues (professionnalisation de notre activité LCD / LMNP)",
    "Pierre-Bénite = première opération commune Marc & Antoine d'une série",
    "Approche structurée et chiffrée — montée en charge maîtrisée",
]
for j, point in enumerate(points):
    tb = s.shapes.add_textbox(Inches(0.9), Inches(4.3 + j * 0.4), Inches(11.5), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "•  " + point
    for r in p.runs:
        r.font.name = FONT_B
        r.font.size = Pt(13)
        r.font.color.rgb = WHITE

txt(s, 0.9, 5.7, 11.6, 0.4, "Contrepartie commerciale", size=14, bold=True, color=GOLD, font=FONT_B)
counterparts = [
    "Rapatriement de nos comptes professionnels (Marc + Antoine + sociétés)",
    "Cartes bancaires pro et perso",
    "Prévoyance, assurances, produits patrimoniaux à terme",
]
for j, point in enumerate(counterparts):
    tb = s.shapes.add_textbox(Inches(0.9), Inches(6.15 + j * 0.35), Inches(11.5), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✓  " + point
    for r in p.runs:
        r.font.name = FONT_B
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE


# ============ SLIDE 13 : MERCI ============
s = blank(); bg(s, ANTHRACITE)
rect(s, 0, 3.7, 13.333, 0.04, color=GOLD)
txt(s, 0.5, 2.3, 12.3, 1, "Merci", size=72, bold=True, color=WHITE, font=FONT_T, align=PP_ALIGN.CENTER)
txt(s, 0.5, 4, 12.3, 0.5, "À votre disposition pour échanger", size=18, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 0.5, 5.8, 12.3, 0.4, "Marc Chenut  &  Antoine Bocquet", size=16, color=GOLD, font=FONT_B, align=PP_ALIGN.CENTER)
txt(s, 0.5, 6.25, 12.3, 0.4, "Annexes : BP financier détaillé (Excel) + note marché (PDF)",
    size=11, color=WHITE, font=FONT_B, align=PP_ALIGN.CENTER)


prs.save(OUT)
print(f"OK: {OUT}")
